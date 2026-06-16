from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs"
OUT = OUT_DIR / "MuseCut_AI情感创作平台_MVP展示.pptx"

SHOWCASE = ROOT / "frontend" / "public" / "showcase"
ASSETS = ROOT / "docs" / "ppt_assets"

FONT = "Microsoft YaHei"
BG = RGBColor(247, 243, 235)
INK = RGBColor(31, 45, 57)
MUTED = RGBColor(100, 115, 128)
TEAL = RGBColor(45, 137, 125)
TEAL_DARK = RGBColor(28, 67, 75)
CORAL = RGBColor(195, 88, 76)
SAND = RGBColor(235, 225, 209)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def add_text(slide, text, x, y, w, h, size=24, bold=False, color=INK, align=PP_ALIGN.LEFT, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_label(slide, text, x=0.65, y=0.38):
    add_text(slide, text.upper(), x, y, 3.2, 0.3, size=13, bold=True, color=TEAL)


def add_title(slide, title, subtitle=None, y=0.75):
    add_text(slide, title, 0.65, y, 8.2, 0.65, size=30, bold=True, color=INK)
    if subtitle:
        add_text(slide, subtitle, 0.68, y + 0.68, 8.7, 0.55, size=14, color=MUTED)


def rounded_rect(slide, x, y, w, h, fill=WHITE, line=rgb("#e5ded3"), radius=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def pill(slide, text, x, y, w, color=TEAL, text_color=WHITE):
    shape = rounded_rect(slide, x, y, w, 0.42, fill=color, line=color)
    add_text(slide, text, x, y + 0.08, w, 0.25, size=12, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    return shape


def add_bullets(slide, items, x, y, w, h, size=17, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
    return box


def add_image(slide, path, x, y, w=None, h=None):
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def add_footer(slide, n):
    add_text(slide, f"MuseCut AI 情感创作平台 · {n:02d}", 10.55, 7.12, 2.15, 0.22, size=8, color=rgb("#9aa3a8"), align=PP_ALIGN.RIGHT)


def make_presentation():
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []

    # 1 Cover
    s = prs.slides.add_slide(blank); set_bg(s, TEAL_DARK); slides.append(s)
    add_image(s, SHOWCASE / "musecut-hero.png", 6.2, 0, w=7.14, h=4.02)
    rounded_rect(s, 0.65, 0.65, 1.05, 1.05, fill=TEAL, line=TEAL)
    add_text(s, "♡", 0.9, 0.83, 0.4, 0.4, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "MuseCut", 0.68, 2.05, 6.0, 0.85, size=42, bold=True, color=WHITE)
    add_text(s, "面向普通用户的 AI 情感创作平台", 0.72, 2.95, 6.1, 0.5, size=20, color=rgb("#d8eee9"))
    add_text(s, "从真实情绪出发，用 AI 完成声音叙事、视频配乐与社区表达", 0.72, 3.62, 5.7, 0.7, size=17, color=rgb("#f3eee5"))
    pill(s, "15 页以内展示版", 0.72, 5.48, 1.9, color=CORAL)
    pill(s, "5 分钟介绍片配套", 2.78, 5.48, 2.25, color=TEAL)

    # 2 Executive summary
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "01 summary"); add_title(s, "项目一句话", "MuseCut 不是单纯生成音乐，而是帮助普通人把真实经历组织成可观看、可聆听、可分享的情感作品。")
    quote = rounded_rect(s, 0.8, 2.05, 11.7, 2.05, fill=WHITE)
    add_text(s, "AI 负责理解素材、编排声音、降低制作门槛；\n故事、记忆和情感始终属于用户自己。", 1.25, 2.45, 10.6, 1.1, size=30, bold=True, color=INK, align=PP_ALIGN.CENTER)
    for i, text in enumerate(["影像", "文字", "声音", "音乐", "社区"]):
        pill(s, text, 2.25 + i * 1.65, 5.15, 1.05, color=[TEAL, CORAL, TEAL_DARK, TEAL, CORAL][i])

    # 3 Background and pain
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "02 problem"); add_title(s, "普通人的创作痛点", "AI 让生成变容易，但“表达清楚”和“情绪连贯”仍然很难。")
    cards = [
        ("素材零散", "视频、照片、旁白和文字各自分散，缺少故事结构。"),
        ("声音设计门槛高", "BGM、鼓点、转场、留白、人声避让需要后期经验。"),
        ("模板不够真诚", "套模板很快，但容易变成统一风格，削弱个人情感。"),
        ("创作容易中断", "退出后找不到草稿，反复操作让普通用户失去耐心。"),
    ]
    for i, (t, b) in enumerate(cards):
        x = 0.75 + (i % 2) * 6.1
        y = 2.0 + (i // 2) * 1.75
        rounded_rect(s, x, y, 5.55, 1.25, fill=WHITE)
        add_text(s, t, x + 0.32, y + 0.22, 2.2, 0.35, size=19, bold=True)
        add_text(s, b, x + 0.32, y + 0.65, 4.75, 0.36, size=13, color=MUTED)

    # 4 Positioning
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "03 positioning"); add_title(s, "产品定位：AI 声音叙事与情感表达平台")
    add_image(s, SHOWCASE / "creative-lab.png", 7.2, 1.28, w=5.35, h=3.75)
    add_bullets(s, [
        "面向恋爱、亲情、友情、成长、校园与旅行等真实故事。",
        "把视频配乐、音乐生成、故事分享、剧本灵感与 AI 创作教学整合为一个平台。",
        "MVP 以“AI 配乐导演”验证核心价值，再扩展到完整故事创作。",
    ], 0.85, 1.72, 5.8, 2.4, size=18)
    pill(s, "不是替代用户表达", 0.9, 5.08, 2.25, color=TEAL_DARK)
    pill(s, "而是帮助用户表达", 3.35, 5.08, 2.25, color=TEAL)

    # 5 MVP workflow
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "04 mvp flow"); add_title(s, "MVP 核心流程", "把复杂后期压缩为 4 步：上传素材 → AI 声音编排与预览 → 导出成片 → 社区分享。")
    steps = [("1", "上传素材", "视频 + 情绪描述"), ("2", "编排与预览", "AI 生成声音时间轴并可试听"), ("3", "导出成片", "音乐叠加到视频生成 MP4"), ("4", "发布社区", "实名/匿名分享与互动")]
    for i, (num, t, b) in enumerate(steps):
        x = 0.6 + i * 3.15
        rounded_rect(s, x, 2.35, 2.75, 2.15, fill=WHITE)
        rounded_rect(s, x + 0.25, 2.72, 0.68, 0.68, fill=TEAL if i < 2 else SAND, line=TEAL if i < 2 else SAND)
        add_text(s, num, x + 0.45, 2.85, 0.25, 0.22, size=18, bold=True, color=WHITE if i < 2 else INK, align=PP_ALIGN.CENTER)
        add_text(s, t, x + 0.25, 3.72, 2.2, 0.35, size=20, bold=True)
        add_text(s, b, x + 0.25, 4.15, 2.15, 0.42, size=12.5, color=MUTED)
        if i < 3:
            add_text(s, "→", x + 2.85, 3.18, 0.35, 0.35, size=26, bold=True, color=TEAL)
    add_text(s, "用户不需要理解专业音轨，AI 先给出完整方案；用户只在关键处微调。", 1.15, 5.55, 11.0, 0.45, size=18, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

    # 6 AI Music Director
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "05 music director"); add_title(s, "核心能力：AI 配乐导演", "关键不是“生成一首歌”，而是判断什么声音应该在什么时间出现。")
    add_image(s, SHOWCASE / "timeline-studio.png", 0.75, 1.62, w=5.25, h=3.68)
    lanes = [("BGM", 6.55, 1.85, 4.7, TEAL), ("Beat", 6.55, 2.72, 3.5, CORAL), ("SFX", 6.55, 3.58, 4.05, rgb("#d2a74f")), ("Ducking", 6.55, 4.45, 2.7, rgb("#6aa58a"))]
    for name, x, y, w, color in lanes:
        add_text(s, name, 6.15, y + 0.05, 0.65, 0.22, size=12, bold=True, color=MUTED)
        rounded_rect(s, x, y, w, 0.38, fill=color, line=color, radius=MSO_SHAPE.RECTANGLE)
    add_bullets(s, ["镜头变化 → 转场与起伏", "字幕重点 → 强调音效", "人声区间 → 自动压低 BGM", "结尾阶段 → 情绪收束"], 6.15, 5.2, 5.7, 1.0, size=15)

    # 7 Sound generation
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "06 audio generation"); add_title(s, "声音生成策略", "MVP 使用 ACE-Step-1.5 生成 BGM，后端负责时长校验、片段适配与视频合成。")
    blocks = [("ACE-Step-1.5", "只允许 ACE 生成 BGM；失败则完整报错，不使用低质量本地兜底。"), ("时长覆盖", "每段音乐按时间轴补齐/裁剪，最终 BGM 覆盖完整视频。"), ("FFmpeg 合成", "导出默认生成带配乐的视频成片 MP4，音频/JSON 作为补充素材。")]
    for i, (t, b) in enumerate(blocks):
        rounded_rect(s, 0.85 + i * 4.05, 2.1, 3.45, 2.25, fill=WHITE)
        add_text(s, t, 1.17 + i * 4.05, 2.48, 2.6, 0.35, size=21, bold=True, color=TEAL_DARK)
        add_text(s, b, 1.17 + i * 4.05, 3.08, 2.7, 0.88, size=13.5, color=MUTED)
    add_text(s, "技术取舍：把壁垒放在“声音编排与作品流程”，而不是绑定某一个音乐模型。", 1.05, 5.5, 11.2, 0.45, size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # 8 Feature matrix
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "07 product modules"); add_title(s, "MVP 功能模块", "围绕普通用户的一次完整创作闭环设计。")
    mods = [("创作入口", "上传视频、表达方向、故事描述"), ("声音编排", "AI 时间轴、点击跳转、带 BGM 试听"), ("预览导出", "MP4 成片、音频、JSON 方案"), ("个人中心", "自动保存、继续任务、删除草稿、头像昵称"), ("故事社区", "动态流、匿名/实名、点赞评论收藏"), ("创作教学", "模板示例、配乐与叙事知识")]
    for i, (t, b) in enumerate(mods):
        x = 0.75 + (i % 3) * 4.15
        y = 1.72 + (i // 3) * 1.85
        rounded_rect(s, x, y, 3.58, 1.38, fill=WHITE)
        add_text(s, t, x + 0.28, y + 0.26, 2.2, 0.3, size=18, bold=True, color=TEAL_DARK)
        add_text(s, b, x + 0.28, y + 0.72, 2.85, 0.36, size=12.8, color=MUTED)

    # 9 Community
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "08 community"); add_title(s, "故事社区：从工具到情感交流平台", "像朋友圈/微博动态一样展示作品，让表达被看见，也让用户保留边界感。")
    add_image(s, SHOWCASE / "community-stories.png", 7.05, 1.35, w=5.35, h=3.75)
    add_bullets(s, [
        "匿名或实名发布作品，降低情绪表达压力。",
        "最近三天展示、不给谁看、不看谁等可见性设置。",
        "支持点赞、评论、收藏、删除作品，形成基础社交闭环。",
        "发布入口放在导出最后一步，避免创作过程中过早打扰。",
    ], 0.85, 1.55, 5.7, 3.25, size=17)

    # 10 Autosave
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "09 user continuity"); add_title(s, "个人中心与自动保存", "用户退出后不需要从头开始，未完成任务可继续、回退或删除。")
    flow = [("自动保存", "创作过程持续写入草稿"), ("个人中心", "头像入口呼出任务与设置"), ("继续创作", "恢复未完成项目状态"), ("删除/管理", "用户控制自己的作品")]
    for i, (t, b) in enumerate(flow):
        rounded_rect(s, 1.0 + i * 3.0, 2.15, 2.45, 1.78, fill=WHITE)
        add_text(s, t, 1.25 + i * 3.0, 2.58, 1.8, 0.32, size=19, bold=True, color=TEAL_DARK)
        add_text(s, b, 1.25 + i * 3.0, 3.05, 1.75, 0.45, size=12.5, color=MUTED)
        if i < 3:
            add_text(s, "→", 3.55 + i * 3.0, 2.72, 0.35, 0.35, size=26, bold=True, color=CORAL)
    add_text(s, "设计目标：把用户的注意力留给表达本身，而不是记住流程和按钮。", 1.1, 5.35, 11.0, 0.42, size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # 11 Architecture
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "10 architecture"); add_title(s, "技术架构", "五层结构，各层通过结构化数据协作，便于替换模型与独立测试。")
    layers = [("前端交互层", "React + Vite + TypeScript"), ("视频分析层", "FFmpeg + 关键点/人声区间"), ("AI 编排层", "LLM 结构化 Music Timeline JSON"), ("声音生成层", "ACE-Step BGM + Beat/SFX/Ducking"), ("导出层", "FFmpeg 混音与 MP4 成片")]
    for i, (t, b) in enumerate(layers):
        y = 1.35 + i * 0.92
        rounded_rect(s, 1.3, y, 10.7, 0.63, fill=WHITE if i % 2 == 0 else rgb("#edf5f2"))
        add_text(s, t, 1.65, y + 0.16, 2.1, 0.22, size=15, bold=True, color=TEAL_DARK)
        add_text(s, b, 4.1, y + 0.16, 5.9, 0.22, size=13.5, color=MUTED)
    pill(s, "标准化 JSON 协议", 5.35, 6.25, 2.55, color=TEAL_DARK)

    # 12 Current status
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "11 mvp status"); add_title(s, "当前完成度", "已完成可运行网页原型，并围绕真实使用反馈持续修复。")
    done = ["注册登录与个人中心", "视频上传与项目管理", "AI 声音编排与 ACE 生成", "带额外音轨的预览试听", "MP4 成片导出", "故事社区与可见性设置", "错误上下文完整提示", "自动保存与任务恢复"]
    for i, item in enumerate(done):
        x = 0.9 + (i % 2) * 5.85
        y = 1.75 + (i // 2) * 0.83
        add_text(s, "✓", x, y, 0.35, 0.25, size=18, bold=True, color=TEAL)
        add_text(s, item, x + 0.38, y + 0.02, 4.6, 0.25, size=15.5, color=INK)
    add_image(s, ASSETS / "export_preview.jpg", 8.8, 4.35, w=3.4, h=1.91)
    add_text(s, "导出示例：视频成片，而非单独音乐", 8.75, 6.35, 3.5, 0.28, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)

    # 13 Demo script
    s = prs.slides.add_slide(blank); set_bg(s); slides.append(s)
    add_label(s, "12 video demo"); add_title(s, "5 分钟介绍片结构", "视频建议以实机录屏为主，PPT 做背景支撑。")
    demo = [("0:00-0:30", "开场：普通人表达真实故事的需求"), ("0:30-1:10", "痛点：找音乐、卡节奏、人声避让、导出复杂"), ("1:10-2:40", "实机：上传视频 → 选择情绪 → AI 生成时间轴"), ("2:40-3:35", "实机：点击时间段跳转、试听额外音轨、微调"), ("3:35-4:25", "实机：导出 MP4 成片 + 发布社区"), ("4:25-5:00", "总结技术架构与平台愿景")]
    for i, (time, desc) in enumerate(demo):
        y = 1.55 + i * 0.76
        pill(s, time, 0.9, y, 1.35, color=TEAL if i < 3 else CORAL)
        add_text(s, desc, 2.55, y + 0.08, 8.7, 0.26, size=15.5, color=INK)

    # 14 Roadmap
    s = prs.slides.add_slide(blank); set_bg(s, TEAL_DARK); slides.append(s)
    add_label(s, "13 roadmap", x=0.65, y=0.45)
    add_text(s, "后续演进", 0.7, 0.95, 5.5, 0.65, size=34, bold=True, color=WHITE)
    add_text(s, "从 AI 配乐导演扩展为完整的 AI 情感故事创作平台", 0.74, 1.65, 7.0, 0.42, size=17, color=rgb("#d8eee9"))
    roadmap = [("MVP", "AI 配乐导演\n短视频声音时间轴"), ("V1", "故事理解引擎\n照片/视频/文字/语音协同"), ("V2", "创作社区\n主题活动与作品沉淀"), ("长期", "导演/编剧灵感库\n真实故事授权改编")]
    for i, (t, b) in enumerate(roadmap):
        rounded_rect(s, 0.9 + i * 3.05, 3.0, 2.55, 1.65, fill=rgb("#f7f3eb"), line=rgb("#f7f3eb"))
        add_text(s, t, 1.2 + i * 3.05, 3.32, 1.2, 0.28, size=18, bold=True, color=TEAL_DARK)
        add_text(s, b, 1.2 + i * 3.05, 3.82, 1.9, 0.55, size=12.8, color=MUTED)
    add_text(s, "让 AI 从“生成素材”走向“辅助表达”，让普通人的真实经历更容易成为作品。", 1.15, 6.1, 11.0, 0.45, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, slide in enumerate(slides, start=1):
        if i != 1:
            add_footer(slide, i)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(make_presentation())
